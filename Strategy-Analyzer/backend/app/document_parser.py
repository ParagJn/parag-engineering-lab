import base64
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


MAX_EXTRACTED_CHARS = 120_000
# Anthropic image size limit is ~5 MB base64; keep raw bytes under 3.5 MB
_MAX_IMAGE_BYTES = 3_500_000
_MIN_IMAGE_BYTES = 128


def _trim(text: str) -> str:
    cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return cleaned[:MAX_EXTRACTED_CHARS]


def _detect_media_type(data: bytes) -> str:
    """Return MIME type based on magic bytes."""
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return "image/png"
    if data[:2] == b"\xff\xd8":
        return "image/jpeg"
    if data[:4] == b"GIF8":
        return "image/gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "image/webp"
    return "image/jpeg"


def _image_entry(data: bytes, location: str) -> dict[str, str] | None:
    """Return a dict suitable for vision processing, or None if image should be skipped."""
    if not (_MIN_IMAGE_BYTES <= len(data) <= _MAX_IMAGE_BYTES):
        return None
    media_type = _detect_media_type(data)
    return {
        "data": base64.b64encode(data).decode("utf-8"),
        "media_type": media_type,
        "location": location,
    }


def extract_images_from_pdf(path: Path) -> list[dict[str, str]]:
    """Extract embedded images from all pages of a PDF."""
    reader = PdfReader(str(path))
    images: list[dict[str, str]] = []
    for page_num, page in enumerate(reader.pages, start=1):
        try:
            for img_obj in page.images:
                entry = _image_entry(img_obj.data, f"Page {page_num}")
                if entry:
                    images.append(entry)
        except Exception:
            pass
    return images


def extract_images_from_docx(path: Path) -> list[dict[str, str]]:
    """Extract embedded images from a DOCX file."""
    doc = Document(str(path))
    images: list[dict[str, str]] = []
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    data = rel.target_part.blob
                    entry = _image_entry(data, "Document")
                    if entry:
                        images.append(entry)
                except Exception:
                    pass
    except Exception:
        pass
    return images


def extract_images_from_pptx(path: Path) -> list[dict[str, str]]:
    """Extract embedded images from all slides of a PPTX file."""
    deck = Presentation(str(path))
    images: list[dict[str, str]] = []
    for slide_num, slide in enumerate(deck.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    data = shape.image.blob
                    entry = _image_entry(data, f"Slide {slide_num}")
                    if entry:
                        images.append(entry)
                except Exception:
                    pass
    return images


def extract_images(path: Path) -> list[dict[str, str]]:
    """Dispatch image extraction based on file type."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_images_from_pdf(path)
    if suffix == ".docx":
        return extract_images_from_docx(path)
    if suffix == ".pptx":
        return extract_images_from_pptx(path)
    return []


def parse_pdf(path: Path) -> dict[str, Any]:
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append({"number": index, "title": f"Page {index}", "text": text.strip()})
    return {"kind": "pdf", "pages": pages, "text": _trim("\n\n".join(p["text"] for p in pages))}


def parse_docx(path: Path) -> dict[str, Any]:
    doc = Document(str(path))
    blocks = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))
    text = _trim("\n".join(blocks))
    return {"kind": "docx", "pages": [{"number": 1, "title": "Document", "text": text}], "text": text}


def parse_pptx(path: Path) -> dict[str, Any]:
    deck = Presentation(str(path))
    slides = []
    for index, slide in enumerate(deck.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                lines.append(shape.text.strip())
        title = lines[0].splitlines()[0][:80] if lines else f"Slide {index}"
        slides.append({"number": index, "title": title, "text": "\n".join(lines)})
    return {"kind": "pptx", "pages": slides, "text": _trim("\n\n".join(s["text"] for s in slides))}


def parse_document(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix == ".docx":
        return parse_docx(path)
    if suffix == ".pptx":
        return parse_pptx(path)
    raise ValueError("Unsupported file type. Upload a .pdf, .docx, or .pptx file.")
