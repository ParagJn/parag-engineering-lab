"""
FastAPI Backend for Project Schedule Tool

This backend provides AI-powered intelligence to the frontend React application.
It uses the Universal LLM Client to interact with various AI models.

Author: Auto-generated
Date: 2026-07-29
"""

from fastapi import FastAPI, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from typing import Optional, Dict, Any, List, Literal
import asyncio
import logging
import json
import os
import re
import shutil
import subprocess
import tempfile
import io
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv

from llm_client import UniversalLLMClient, LLMClientError
from ibm_ica_client import IBMICAClient, IBMICAError
from pptx import Presentation
from pptx.util import Inches, Pt

# Load environment variables
load_dotenv()

# Initialize FastAPI app
app = FastAPI(
    title="Project Schedule Tool API",
    description="AI-powered backend for project scheduling and planning",
    version="1.0.0"
)

# Configure CORS for frontend communication
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:3000"],  # Vite and React dev servers
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("ProjectScheduleAPI")

# Global LLM client instances
llm_client: Optional[UniversalLLMClient] = None
ibm_ica_client: Optional[IBMICAClient] = None

# Settings file path
SETTINGS_FILE = Path("settings.json")

# Default settings
DEFAULT_SETTINGS = {
    "ai_provider": "sap"  # "sap" or "ibm_ica"
}

def load_settings() -> Dict[str, Any]:
    """Load settings from file or return defaults"""
    if SETTINGS_FILE.exists():
        try:
            with open(SETTINGS_FILE, 'r') as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"Failed to load settings: {e}")
    return DEFAULT_SETTINGS.copy()

def save_settings(settings: Dict[str, Any]) -> bool:
    """Save settings to file"""
    try:
        with open(SETTINGS_FILE, 'w') as f:
            json.dump(settings, f, indent=2)
        return True
    except Exception as e:
        logger.error(f"Failed to save settings: {e}")
        return False

def get_active_provider() -> str:
    """Get the currently selected AI provider"""
    settings = load_settings()
    return settings.get("ai_provider", "sap")


# Pydantic Models for Request/Response
class LLMRequest(BaseModel):
    """Request model for LLM generation"""
    prompt: str = Field(..., description="The prompt to send to the LLM")
    model: Optional[str] = Field(None, description="Specific model to use (leave empty for default)")
    provider: Optional[str] = Field(None, description="Provider to use (leave empty for default)")
    max_tokens: Optional[int] = Field(25000, description="Maximum tokens to generate")


class LLMResponse(BaseModel):
    """Response model for LLM generation"""
    success: bool
    response: Optional[Dict[Any, Any]] = None
    error: Optional[str] = None
    model_used: Optional[str] = None
    provider_used: Optional[str] = None
    timestamp: str


class ModelInfo(BaseModel):
    """Model information"""
    provider: str
    model_name: str
    capabilities: List[str]
    is_default: bool
    max_tokens: Optional[int]


class HealthResponse(BaseModel):
    """Health check response"""
    status: str
    timestamp: str
    llm_client_initialized: bool


class TaskAnalysisRequest(BaseModel):
    """Request for task analysis"""
    task_description: str
    context: Optional[Dict[str, Any]] = None


class TaskAnalysisResponse(BaseModel):
    """Response for task analysis"""
    success: bool
    analysis: Optional[Dict[str, Any]] = None
    error: Optional[str] = None


class ProjectOptimizationRequest(BaseModel):
    """Request for project optimization"""
    project_data: Dict[str, Any]
    optimization_goals: List[str]


class ProjectOptimizationResponse(BaseModel):
    """Response for project optimization"""
    success: bool
    optimizations: Optional[List[Dict[str, Any]]] = None
    error: Optional[str] = None


class SoWGenerationRequest(BaseModel):
    """Request for Statement of Work generation"""
    project_name: str
    customer: str
    background: Optional[str] = None
    assumptions: Optional[str] = None
    out_of_scope: Optional[str] = None
    maw_deliverables: Optional[str] = None
    tasks: Optional[List[Dict[str, Any]]] = None


class SoWGenerationResponse(BaseModel):
    """Response for Statement of Work generation"""
    success: bool
    sow_content: Optional[str] = None
    needs_more_info: bool = False
    questions: Optional[List[str]] = None
    error: Optional[str] = None
    timestamp: str


# Dependency to get LLM client
def get_llm_client() -> UniversalLLMClient:
    """Dependency to get the LLM client instance"""
    global llm_client
    if llm_client is None:
        raise HTTPException(status_code=503, detail="LLM client not initialized")
    return llm_client


# Startup event
@app.on_event("startup")
async def startup_event():
    """Initialize services on startup"""
    global llm_client, ibm_ica_client
    
    # Initialize SAP AI Core client
    try:
        llm_client = UniversalLLMClient(config_path="config.json")
        logger.info("✅ SAP AI Core LLM Client initialized successfully")
        logger.info(f"Available models: {len(llm_client.list_available_models())}")
    except Exception as e:
        logger.error(f"❌ Failed to initialize SAP LLM client: {e}")
        llm_client = None
    
    # Initialize IBM ICA client
    try:
        ibm_endpoint = os.getenv("IBM_ICA_ENDPOINT")
        ibm_api_key = os.getenv("IBM_ICA_API_KEY")
        ibm_model_id = os.getenv("IBM_ICA_MODEL_ID", "claude-sonnet-5")
        
        if ibm_endpoint and ibm_api_key:
            ibm_ica_client = IBMICAClient(
                endpoint=ibm_endpoint,
                api_key=ibm_api_key,
                model_id=ibm_model_id
            )
            logger.info("✅ IBM ICA Client initialized successfully")
            logger.info(f"IBM ICA Model: {ibm_model_id}")
        else:
            logger.warning("⚠️  IBM ICA credentials not found in environment")
            ibm_ica_client = None
    except Exception as e:
        logger.error(f"❌ Failed to initialize IBM ICA client: {e}")
        ibm_ica_client = None
    
    # Log active provider
    active_provider = get_active_provider()
    logger.info(f"Active AI Provider: {active_provider.upper()}")


# Shutdown event
@app.on_event("shutdown")
async def shutdown_event():
    """Cleanup on shutdown"""
    logger.info("Shutting down API server")


# Health check endpoint
@app.get("/health", response_model=HealthResponse)
async def health_check():
    """Health check endpoint"""
    return HealthResponse(
        status="healthy",
        timestamp=datetime.now().isoformat(),
        llm_client_initialized=llm_client is not None
    )


# Root endpoint
@app.get("/")
async def root():
    """Root endpoint"""
    return {
        "message": "Project Schedule Tool API",
        "version": "1.0.0",
        "docs": "/docs",
        "health": "/health"
    }


# Settings endpoints
class ProviderSettingsRequest(BaseModel):
    """Request to update provider settings"""
    ai_provider: str = Field(..., description="AI provider: 'sap' or 'ibm_ica'")


class ProviderSettingsResponse(BaseModel):
    """Response for provider settings"""
    success: bool
    ai_provider: str
    sap_available: bool
    ibm_ica_available: bool
    message: Optional[str] = None


@app.get("/settings/provider", response_model=ProviderSettingsResponse)
async def get_provider_settings():
    """Get current provider settings"""
    settings = load_settings()
    return ProviderSettingsResponse(
        success=True,
        ai_provider=settings.get("ai_provider", "sap"),
        sap_available=llm_client is not None,
        ibm_ica_available=ibm_ica_client is not None
    )


@app.post("/settings/provider", response_model=ProviderSettingsResponse)
async def update_provider_settings(request: ProviderSettingsRequest):
    """Update provider settings"""
    if request.ai_provider not in ["sap", "ibm_ica"]:
        raise HTTPException(status_code=400, detail="Invalid provider. Must be 'sap' or 'ibm_ica'")
    
    # Check if the requested provider is available
    if request.ai_provider == "sap" and llm_client is None:
        raise HTTPException(status_code=503, detail="SAP AI Core client not available")
    
    if request.ai_provider == "ibm_ica" and ibm_ica_client is None:
        raise HTTPException(status_code=503, detail="IBM ICA client not available")
    
    settings = load_settings()
    settings["ai_provider"] = request.ai_provider
    
    if save_settings(settings):
        logger.info(f"Provider switched to: {request.ai_provider.upper()}")
        return ProviderSettingsResponse(
            success=True,
            ai_provider=request.ai_provider,
            sap_available=llm_client is not None,
            ibm_ica_available=ibm_ica_client is not None,
            message=f"Provider successfully switched to {request.ai_provider.upper()}"
        )
    else:
        raise HTTPException(status_code=500, detail="Failed to save settings")


# List available models
@app.get("/models", response_model=List[ModelInfo])
async def list_models(client: UniversalLLMClient = Depends(get_llm_client)):
    """List all available AI models"""
    try:
        models = client.list_available_models()
        return [ModelInfo(**model) for model in models]
    except Exception as e:
        logger.error(f"Error listing models: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# Get default model
@app.get("/models/default")
async def get_default_model(client: UniversalLLMClient = Depends(get_llm_client)):
    """Get the default model configuration"""
    try:
        provider, model = client.get_default_model()
        return {
            "provider": provider,
            "model": model
        }
    except Exception as e:
        logger.error(f"Error getting default model: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# Generic LLM generation endpoint
@app.post("/generate", response_model=LLMResponse)
async def generate_text(
    request: LLMRequest,
    client: UniversalLLMClient = Depends(get_llm_client)
):
    """
    Generate text using the LLM.
    
    This is a generic endpoint that can be used for any text generation task.
    """
    # Get model info upfront for error reporting
    provider_name = request.provider if request.provider else client.config['settings']['default_provider']
    model_name = request.model if request.model else client.config['settings']['default_model']
    
    try:
        # Prepare kwargs
        kwargs = {}
        if request.max_tokens:
            kwargs['max_tokens'] = request.max_tokens
        
        # Generate response
        response = client.generate(
            prompt=request.prompt,
            model=request.model if request.model else None,
            provider=request.provider if request.provider else None,
            **kwargs
        )
        
        return LLMResponse(
            success=True,
            response=response,
            model_used=model_name,
            provider_used=provider_name,
            timestamp=datetime.now().isoformat()
        )
        
    except LLMClientError as e:
        logger.error(f"LLM Client error: {e}")
        return LLMResponse(
            success=False,
            error=str(e),
            model_used=model_name,
            provider_used=provider_name,
            timestamp=datetime.now().isoformat()
        )
    except Exception as e:
        logger.error(f"Unexpected error: {e}")
        return LLMResponse(
            success=False,
            error=f"Internal server error: {str(e)}",
            model_used=model_name,
            provider_used=provider_name,
            timestamp=datetime.now().isoformat()
        )


# Task analysis endpoint
@app.post("/analyze/task", response_model=TaskAnalysisResponse)
async def analyze_task(
    request: TaskAnalysisRequest,
    client: UniversalLLMClient = Depends(get_llm_client)
):
    """
    Analyze a task description and provide structured insights.
    
    This endpoint uses AI to break down tasks, estimate effort,
    identify dependencies, and suggest optimal scheduling.
    """
    try:
        # Build analysis prompt
        context_str = ""
        if request.context:
            context_str = f"\n\nContext:\n{request.context}"
        
        prompt = f"""Analyze the following task and provide a structured breakdown:

Task: {request.task_description}{context_str}

Please provide:
1. A refined task title (concise, action-oriented)
2. Estimated effort in person-days (provide a range: min, max, likely)
3. Key dependencies or prerequisites
4. Potential risks or challenges
5. Suggested sub-tasks if the task is complex
6. Priority recommendation (High, Medium, Low) with justification

Format your response as JSON with the following structure:
{{
    "title": "string",
    "effort_estimate": {{"min": number, "max": number, "likely": number}},
    "dependencies": ["string"],
    "risks": ["string"],
    "sub_tasks": ["string"],
    "priority": {{"level": "string", "justification": "string"}}
}}"""
        
        response = client.generate(prompt=prompt)
        
        # Parse the response (this is simplified - you might need more robust parsing)
        return TaskAnalysisResponse(
            success=True,
            analysis=response
        )
        
    except Exception as e:
        logger.error(f"Task analysis error: {e}")
        return TaskAnalysisResponse(
            success=False,
            error=str(e)
        )


# Project optimization endpoint
@app.post("/optimize/project", response_model=ProjectOptimizationResponse)
async def optimize_project(
    request: ProjectOptimizationRequest,
    client: UniversalLLMClient = Depends(get_llm_client)
):
    """
    Optimize a project schedule based on given goals.
    
    This endpoint uses AI to suggest optimizations for:
    - Resource allocation
    - Task sequencing
    - Parallel execution opportunities
    - Critical path optimization
    """
    try:
        prompt = f"""Analyze this project and suggest optimizations:

Project Data:
{request.project_data}

Optimization Goals:
{', '.join(request.optimization_goals)}

Provide specific, actionable recommendations for:
1. Resource reallocation opportunities
2. Tasks that can be parallelized
3. Critical path optimizations
4. Risk mitigation strategies
5. Timeline compression opportunities

Format as JSON array of optimization suggestions, each with:
{{
    "category": "string",
    "description": "string",
    "impact": "High|Medium|Low",
    "effort": "Easy|Moderate|Difficult",
    "timeline_savings": "number (in days)"
}}"""
        
        response = client.generate(prompt=prompt, max_tokens=25000)
        
        return ProjectOptimizationResponse(
            success=True,
            optimizations=response
        )
        
    except Exception as e:
        logger.error(f"Project optimization error: {e}")
        return ProjectOptimizationResponse(
            success=False,
            error=str(e)
        )


# Natural language query endpoint
@app.post("/query")
async def natural_language_query(
    query: str,
    context: Optional[Dict[str, Any]] = None,
    client: UniversalLLMClient = Depends(get_llm_client)
):
    """
    Answer natural language queries about project planning.
    
    Examples:
    - "What's the critical path for this project?"
    - "How can I reduce the timeline by 2 weeks?"
    - "What are the risks in the current schedule?"
    """
    try:
        context_str = ""
        if context:
            context_str = f"\n\nProject Context:\n{context}"
        
        prompt = f"""You are an expert project management assistant. Answer the following question:

Question: {query}{context_str}

Provide a clear, concise, and actionable answer based on project management best practices."""
        
        response = client.generate(prompt=prompt)
        
        return {
            "success": True,
            "query": query,
            "answer": response,
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f"Query error: {e}")
        return {
            "success": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }


# Provider wrapper for AI generation
def generate_with_active_provider(prompt: str, max_tokens: int = 25000) -> str:
    """
    Generate content using the currently selected provider.
    
    Returns the generated text content.
    Raises HTTPException if provider is not available.
    """
    active_provider = get_active_provider()
    
    if active_provider == "ibm_ica":
        if ibm_ica_client is None:
            raise HTTPException(
                status_code=503, 
                detail="IBM ICA client not initialized. Check environment variables."
            )
        
        try:
            logger.info("Generating content with IBM ICA...")
            messages = [{"role": "user", "content": prompt}]
            result = ibm_ica_client.chat(messages=messages, max_tokens=max_tokens)
            logger.info(f"IBM ICA generation completed. Tokens: {result.get('total_tokens', 'unknown')}")
            return result.get("text", "")
        except IBMICAError as e:
            logger.error(f"IBM ICA generation failed: {e}")
            raise HTTPException(status_code=500, detail=f"IBM ICA generation failed: {str(e)}")
    
    else:  # default to SAP
        if llm_client is None:
            raise HTTPException(
                status_code=503, 
                detail="SAP AI Core client not initialized. Check configuration."
            )
        
        try:
            logger.info("Generating content with SAP AI Core...")
            response = llm_client.generate(prompt=prompt, max_tokens=max_tokens)
            
            # Extract content from SAP AI Core orchestration response
            content = ""
            
            if isinstance(response, dict):
                # SAP AI Core returns: final_result.choices[0].message.content
                if 'final_result' in response:
                    final = response['final_result']
                    if isinstance(final, dict) and 'choices' in final:
                        choices = final['choices']
                        if isinstance(choices, list) and len(choices) > 0:
                            choice = choices[0]
                            if isinstance(choice, dict) and 'message' in choice:
                                message = choice['message']
                                if isinstance(message, dict):
                                    content = message.get('content', '')
                                    if content:
                                        logger.info("✓ Extracted content from final_result.choices[0].message.content")
                
                # Fallback: try orchestration_result.choices[0].message.content
                if not content and 'orchestration_result' in response:
                    orch = response['orchestration_result']
                    if isinstance(orch, dict) and 'choices' in orch:
                        choices = orch['choices']
                        if isinstance(choices, list) and len(choices) > 0:
                            choice = choices[0]
                            if isinstance(choice, dict) and 'message' in choice:
                                message = choice['message']
                                if isinstance(message, dict):
                                    content = message.get('content', '')
                
                # Last fallback: try content field directly
                if not content:
                    content = response.get('content', '')
            
            if not content:
                logger.warning("Could not extract content from SAP response")
                content = str(response)
            
            logger.info(f"SAP AI Core generation completed. Length: {len(content)}")
            return content
            
        except LLMClientError as e:
            logger.error(f"SAP generation failed: {e}")
            raise HTTPException(status_code=500, detail=f"SAP generation failed: {str(e)}")


# --- SoW generation: shared grounding + section prompts -------------------

SOW_RULES_PREAMBLE = """You are an experienced business analyst drafting ONE section of a real \
Statement of Work (SoW) for a paying client. Write like a human consultant editing their own \
draft, not a template being auto-filled.

Hard rules:
- Only state facts that are directly supported by the information given below (the background, \
assumptions, out-of-scope notes, deliverables notes, and the project schedule). Do not invent \
specific numbers, counts, system names, dates, or details that are not given or directly \
derivable from the schedule.
- If the given information is genuinely silent on something, leave it out rather than writing \
generic filler to fill space.
- Do not pad length. Write only as much as the material actually supports - a short, accurate \
section beats a long, generic one.
- Vary sentence structure and word choice across bullets and paragraphs. Do not mechanically \
repeat the same stock phrases (e.g. "leverage", "robust", "seamless", or "Key Outcomes" as a \
heading on every block).
- Format in Markdown only: use ## and ### headers, **bold**, bullet points, and | pipe | tables | \
where a table is requested below. No preamble like "Here is the section" - output the section \
content directly, starting with its heading.
"""

SOW_SECTION_TASKS: List[tuple] = [
    (
        "Background and Context",
        """SECTION TASK: Write the "## Background and Context" section.
- 1 to 4 paragraphs (only as many as the material supports) explaining the business context, \
current state, and why this project is needed.
- Base it strictly on the BACKGROUND text given below - do not introduce scope or deliverable \
details here.""",
    ),
    (
        "Scope and Out of Scope",
        """SECTION TASK: Write the "## Scope" and "## Out of Scope" sections.

For Scope:
- Break the work into numbered subsections (### 1.1, ### 1.2, etc.) that map to logical groupings \
of the PROJECT SCHEDULE tasks below (group related tasks together; do not invent scope areas that \
have no corresponding task).
- Each subsection: a short "Scope Overview" paragraph, then a "Scope of Work" bullet list of what \
will actually be done (grounded in the matching tasks). Only add a "Key Outcomes" bullet list if \
there is something concrete to state - skip it otherwise.

For Out of Scope:
- Present as a markdown table: | Scope Area | Details |
- Base rows strictly on the OUT OF SCOPE text below. If nothing was provided, write a single-row \
table noting that exclusions were not specified and should be confirmed with the client - do not \
invent a full exclusions list from scratch.""",
    ),
    (
        "Deliverables and Work Products",
        """SECTION TASK: Write the "## Deliverables and Work Products" section.
- "### Deliverables" table: | Item | Description | Target Date/Phase | - formal outputs given to \
the client. Map these to the PROJECT SCHEDULE tasks and their dates/phases where possible.
- "### Work Products" table: | Item | Description | Purpose | - internal/supporting artifacts \
produced during execution.
- Base both tables primarily on the MAW DELIVERABLES text and the PROJECT SCHEDULE below. If MAW \
DELIVERABLES is "Not provided", derive the list directly from the schedule's tasks only - do not \
invent deliverables unrelated to any task.""",
    ),
    (
        "Risks and Dependencies",
        """SECTION TASK: Write the "## Risks" and "## Dependencies" sections.

For Risks:
- Table: | Risk Category | Description | Impact | Mitigation Strategy |
- Base risks on what the BACKGROUND, ASSUMPTIONS, and PROJECT SCHEDULE actually suggest (e.g. \
tight dependency chains, tasks with no buffer, external dependencies mentioned in the text). \
Impact is High/Medium/Low. Do not pad this out with a fixed catalog of generic risk categories - \
include only risks the input actually supports.

For Dependencies:
- Table: | Dependency Type | Description | Owner | Status |
- Base rows on dependencies stated in ASSUMPTIONS/BACKGROUND and on task-level dependency \
relationships in the PROJECT SCHEDULE. Owner is IBM, Client, or Third-party; Status is Required, \
Optional, or Critical.""",
    ),
    (
        "RACI Matrix",
        """SECTION TASK: Write the "## RACI Matrix" section.
- Table: | Activity/Task | IBM | Client |
- Use the PROJECT SCHEDULE task list below as the row source - one row per task/activity (group \
very granular sub-activities under their parent task if there are many). Do not invent activities \
that aren't in the schedule.
- Use standard RACI letters (R/A/C/I), multiple letters per cell where appropriate.""",
    ),
]


def _render_task_list(tasks: Optional[List[Dict[str, Any]]]) -> str:
    """Render the Planner's task list into a compact, factual plain-text block."""
    if not tasks:
        return "Not provided."

    lines: List[str] = []
    for t in tasks:
        index = t.get("index", "?")
        activity = t.get("activity", "Untitled task")
        days = t.get("estimatedDays", "?")
        weeks = t.get("estimatedWeeks", "?")
        fte = t.get("fte", "?")
        dependency = t.get("dependency") or "None"
        start = t.get("calculatedStartDate") or "TBD"
        finish = t.get("calculatedFinishDate") or "TBD"
        lines.append(
            f"{index}. {activity} - {days} day(s), {weeks} week(s), FTE {fte}, "
            f"depends on: {dependency}, {start} to {finish}"
        )
        for sub in (t.get("subActivities") or []):
            lines.append(f"    - {sub}")

    return "\n".join(lines)


def _build_shared_sow_context(request: "SoWGenerationRequest") -> str:
    """Build the grounding context block shared by every SoW section prompt."""
    return f"""PROJECT INFORMATION:
- Project Name: {request.project_name}
- Customer: {request.customer}

BACKGROUND:
{request.background}

ASSUMPTIONS:
{request.assumptions if request.assumptions else "Not provided"}

OUT OF SCOPE (as described by the user):
{request.out_of_scope if request.out_of_scope else "Not provided"}

MAW DELIVERABLES (as described by the user):
{request.maw_deliverables if request.maw_deliverables else "Not provided"}

PROJECT SCHEDULE (real tasks from the project plan - treat this as ground truth for scope, \
deliverables, and RACI; do not invent tasks beyond this list):
{_render_task_list(request.tasks)}
"""


# Statement of Work generation endpoint
@app.post("/generate/sow", response_model=SoWGenerationResponse)
async def generate_sow(
    request: SoWGenerationRequest
):
    """
    Generate a professional Statement of Work (SoW) document.

    Generates the document as independent sections (Background, Scope & Out of
    Scope, Deliverables, Risks & Dependencies, RACI), each grounded in the same
    background/assumptions/out-of-scope/deliverables text plus the real Planner
    task schedule. Sections are generated concurrently (one AI call per section,
    dispatched in parallel) and stitched together in a fixed order.
    """
    try:
        # Check if we have at least background information
        if not request.background or len(request.background.strip()) < 50:
            return SoWGenerationResponse(
                success=False,
                needs_more_info=True,
                questions=[
                    "Please provide detailed project background information (at least a few paragraphs).",
                    "What is the business context and problem this project aims to solve?",
                    "What are the key objectives and expected outcomes?"
                ],
                timestamp=datetime.now().isoformat()
            )

        shared_context = _build_shared_sow_context(request)

        section_prompts = [
            (name, SOW_RULES_PREAMBLE + "\n" + shared_context + "\n" + task)
            for name, task in SOW_SECTION_TASKS
        ]

        logger.info(f"Generating {len(section_prompts)} SoW sections concurrently...")
        loop = asyncio.get_running_loop()
        results = await asyncio.gather(*[
            loop.run_in_executor(None, generate_with_active_provider, prompt, 6000)
            for _, prompt in section_prompts
        ])

        cleaned_sections = []
        for (name, _), content in zip(section_prompts, results):
            content = content.strip()
            # Replace escaped newlines with actual newlines if they exist as literal strings
            if '\\n' in content:
                content = content.replace('\\n\\n', '\n\n').replace('\\n', '\n')
            if content:
                cleaned_sections.append(content)
            else:
                logger.warning(f"SoW section '{name}' returned empty content")

        content = "\n\n".join(cleaned_sections).strip()

        logger.info(f"Final content length: {len(content)} chars")

        # Final validation
        if not content or len(content) < 50:
            return SoWGenerationResponse(
                success=False,
                error="Generated content is too short or empty. Please try again.",
                needs_more_info=False,
                timestamp=datetime.now().isoformat()
            )

        return SoWGenerationResponse(
            success=True,
            sow_content=content,
            needs_more_info=False,
            timestamp=datetime.now().isoformat()
        )

    except Exception as e:
        logger.error(f"SoW generation error: {e}")
        return SoWGenerationResponse(
            success=False,
            error=str(e),
            needs_more_info=False,
            timestamp=datetime.now().isoformat()
        )


# --- ISBD slide deck generation: grounded in the already-generated SoW ----

class ISBDGenerationRequest(BaseModel):
    """Request for ISBD slide deck generation"""
    project_name: str
    customer: str
    sow_content: str


ISBD_JSON_RULES = """You are condensing ONE part of an already-approved Statement of Work (SoW) into \
talking points for a short slide deck. Write like a consultant preparing their own slide notes, not \
a template being auto-filled.

Hard rules:
- Only use facts that are already stated in the SoW text given below. Do not invent numbers, names, \
dates, or details that are not present in it.
- Where the SoW is silent on something, return fewer items (or an empty list) rather than padding \
with generic filler.
- Keep every item short - a slide bullet, not a paragraph. Quality over quantity: a short, accurate \
list beats a long, generic one.
- Output ONLY raw JSON matching the exact shape requested below. No markdown, no code fences, no \
commentary before or after it.
"""


def _extract_json(content: str) -> Optional[Dict[str, Any]]:
    """Best-effort extraction of a single JSON object from a model response."""
    content = content.strip()
    start_idx = content.find('{')
    end_idx = content.rfind('}') + 1
    if start_idx == -1 or end_idx <= start_idx:
        return None
    try:
        return json.loads(content[start_idx:end_idx])
    except json.JSONDecodeError:
        return None


def _build_isbd_prompts(sow_content: str) -> Dict[str, str]:
    """Build the three ISBD section prompts, each grounded in the full SoW text."""
    sow_block = f"STATEMENT OF WORK (ground truth - condense from this only):\n---\n{sow_content}\n---\n"

    approach_prompt = ISBD_JSON_RULES + "\n" + sow_block + """
SECTION TASK: Produce JSON with this exact shape:
{"approach_bullets": ["...", ...], "scope_bullets": ["...", ...], "mermaid": "flowchart LR\\n..."}

- "approach_bullets": up to 5 short bullets summarizing HOW the work will be delivered, drawn from \
the SoW's Background/Scope sections.
- "scope_bullets": up to 6 short bullets summarizing WHAT is in scope, drawn from the SoW's Scope \
section (its numbered subsections/tasks).
- "mermaid": a single minimal Mermaid flowchart definition, "flowchart LR" style, with 3 to 6 nodes \
representing the high-level phases of the approach (2-4 words per node, e.g. Discovery, Design, \
Build, Test, Deploy) IN THE ORDER they occur, connected with simple arrows (A --> B). Base the \
phases strictly on the SoW's Scope subsections - do not invent phases that aren't implied by it. \
Use only letters, digits, spaces, and hyphens inside node labels - no quotes, no special characters, \
no parentheses.
"""

    assumptions_prompt = ISBD_JSON_RULES + "\n" + sow_block + """
SECTION TASK: Produce JSON with this exact shape:
{"assumptions": ["...", ...], "dependencies": ["...", ...]}

- "assumptions": at most 5 short bullets, drawn from the SoW's Assumptions content. Fewer than 5 is \
fine if that's all the SoW supports.
- "dependencies": at most 5 short bullets, drawn from the SoW's Dependencies content. Fewer than 5 \
is fine if that's all the SoW supports.
"""

    risks_prompt = ISBD_JSON_RULES + "\n" + sow_block + """
SECTION TASK: Produce JSON with this exact shape:
{"risks": [{"risk": "...", "impact": "High|Medium|Low", "mitigation": "..."}, ...]}

- At most 5 rows, drawn from the SoW's Risks table/content. Fewer than 5 is fine if that's all the \
SoW supports. Each field should be a short phrase, not a paragraph.
"""

    return {
        "approach": approach_prompt,
        "assumptions": assumptions_prompt,
        "risks": risks_prompt,
    }


def _sanitize_mermaid(mermaid_text: Optional[str]) -> Optional[str]:
    """Strip markdown code-fence wrapping a model may add around a mermaid definition."""
    if not mermaid_text or not mermaid_text.strip():
        return None
    text = mermaid_text.strip()
    text = re.sub(r"^```(?:mermaid)?\s*", "", text)
    text = re.sub(r"```\s*$", "", text)
    text = text.strip()
    return text or None


def _find_mmdc_binary() -> Optional[str]:
    path = shutil.which("mmdc")
    if path:
        return path
    for candidate in ("/opt/homebrew/bin/mmdc", "/usr/local/bin/mmdc"):
        if os.path.isfile(candidate):
            return candidate
    return None


def _render_mermaid_to_png(mermaid_text: str) -> Optional[bytes]:
    """Render a Mermaid diagram definition to a PNG locally via mmdc (mermaid-cli). Returns None on any failure."""
    mmdc_path = _find_mmdc_binary()
    if not mmdc_path:
        logger.warning("mmdc (mermaid-cli) not found on PATH - skipping ISBD diagram")
        return None

    try:
        with tempfile.TemporaryDirectory() as tmp_dir:
            in_path = os.path.join(tmp_dir, "diagram.mmd")
            out_path = os.path.join(tmp_dir, "diagram.png")
            with open(in_path, "w", encoding="utf-8") as f:
                f.write(mermaid_text)

            subprocess.run(
                [mmdc_path, "-i", in_path, "-o", out_path, "-b", "white", "-w", "1400", "-H", "450"],
                capture_output=True,
                timeout=30,
                check=True,
            )

            with open(out_path, "rb") as f:
                return f.read()
    except Exception as e:
        logger.warning(f"Mermaid rendering failed, skipping ISBD diagram: {e}")
        return None


def _set_run_font(run, size: int = 12, bold: bool = False):
    """Apply the ISBD deck's shared font styling to a text run."""
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold


def _build_isbd_presentation(
    approach_bullets: List[str],
    scope_bullets: List[str],
    diagram_png: Optional[bytes],
    assumptions: List[str],
    dependencies: List[str],
    risks: List[Dict[str, str]],
) -> io.BytesIO:
    """Build the 5-slide ISBD deck and return it as an in-memory .pptx stream."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    margin = Inches(0.5)
    content_w = slide_w - 2 * margin

    def add_title(slide, text: str):
        box = slide.shapes.add_textbox(margin, Inches(0.3), content_w, Inches(0.6))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _set_run_font(run, size=14, bold=True)
        return box

    def add_subheader(slide, text: str, left, top, width):
        box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _set_run_font(run, size=14, bold=True)
        return box

    def add_bullets(slide, items: List[str], left, top, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {item}"
            _set_run_font(run, size=12, bold=False)
        return box

    # Slide 1: Approach and Scope
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide1, "Approach and Scope")

    left_col_w = Inches(4.6)
    add_subheader(slide1, "Approach", margin, Inches(1.1), left_col_w)
    add_bullets(slide1, approach_bullets, margin, Inches(1.55), left_col_w, Inches(2.6))

    add_subheader(slide1, "Scope", margin, Inches(4.3), left_col_w)
    add_bullets(slide1, scope_bullets, margin, Inches(4.75), left_col_w, Inches(2.3))

    if diagram_png:
        img_stream = io.BytesIO(diagram_png)
        diagram_left = margin + left_col_w + Inches(0.3)
        diagram_w = slide_w - diagram_left - margin
        slide1.shapes.add_picture(img_stream, diagram_left, Inches(2.2), width=diagram_w)

    # Slides 2 & 3: Schedule / Cost - title only, native "click to add text" placeholder left untouched
    for title_text in ("Schedule", "Cost"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_ph = slide.shapes.title
        title_ph.text_frame.text = title_text
        run = title_ph.text_frame.paragraphs[0].runs[0]
        _set_run_font(run, size=14, bold=True)

    # Slide 4: Assumptions and Dependencies
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide4, "Assumptions and Dependencies")

    col_w = (content_w - Inches(0.4)) // 2
    add_subheader(slide4, "Assumptions", margin, Inches(1.1), col_w)
    add_bullets(slide4, assumptions[:5], margin, Inches(1.55), col_w, Inches(5.3))

    right_col_left = margin + col_w + Inches(0.4)
    add_subheader(slide4, "Dependencies", right_col_left, Inches(1.1), col_w)
    add_bullets(slide4, dependencies[:5], right_col_left, Inches(1.55), col_w, Inches(5.3))

    # Slide 5: Risks table
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide5, "Risks")

    risk_rows = risks[:5]
    rows = len(risk_rows) + 1
    table_shape = slide5.shapes.add_table(rows, 3, margin, Inches(1.2), content_w, Inches(0.5) * rows)
    table = table_shape.table
    table.columns[0].width = int(content_w * 0.45)
    table.columns[1].width = int(content_w * 0.15)
    table.columns[2].width = int(content_w * 0.40)

    headers = ["Risk", "Impact", "Mitigation"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        run = cell.text_frame.paragraphs[0].runs[0]
        _set_run_font(run, size=14, bold=True)

    for row_idx, risk in enumerate(risk_rows, start=1):
        values = [risk.get("risk", ""), risk.get("impact", ""), risk.get("mitigation", "")]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            run = cell.text_frame.paragraphs[0].runs[0]
            _set_run_font(run, size=12, bold=False)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


@app.post("/generate/isbd")
async def generate_isbd(request: ISBDGenerationRequest):
    """
    Generate a 5-slide ISBD PowerPoint deck (Approach & Scope, Schedule, Cost,
    Assumptions & Dependencies, Risks) grounded in the already-generated SoW
    draft. Schedule and Cost slides are left blank for the user to paste into.
    """
    if not request.sow_content or len(request.sow_content.strip()) < 50:
        raise HTTPException(status_code=400, detail="A generated SoW draft is required before creating ISBD slides.")

    try:
        prompts = _build_isbd_prompts(request.sow_content)
        loop = asyncio.get_running_loop()
        approach_raw, assumptions_raw, risks_raw = await asyncio.gather(*[
            loop.run_in_executor(None, generate_with_active_provider, prompts["approach"], 2000),
            loop.run_in_executor(None, generate_with_active_provider, prompts["assumptions"], 1500),
            loop.run_in_executor(None, generate_with_active_provider, prompts["risks"], 1500),
        ])

        approach_json = _extract_json(approach_raw) or {}
        assumptions_json = _extract_json(assumptions_raw) or {}
        risks_json = _extract_json(risks_raw) or {}

        approach_bullets = [str(b) for b in (approach_json.get("approach_bullets") or [])][:6]
        scope_bullets = [str(b) for b in (approach_json.get("scope_bullets") or [])][:6]
        assumptions = [str(a) for a in (assumptions_json.get("assumptions") or [])][:5]
        dependencies = [str(d) for d in (assumptions_json.get("dependencies") or [])][:5]
        raw_risks = risks_json.get("risks") or []
        risks = [r for r in raw_risks if isinstance(r, dict)][:5]

        mermaid_text = _sanitize_mermaid(approach_json.get("mermaid"))
        diagram_png = None
        if mermaid_text:
            diagram_png = await loop.run_in_executor(None, _render_mermaid_to_png, mermaid_text)

        pptx_stream = _build_isbd_presentation(
            approach_bullets=approach_bullets,
            scope_bullets=scope_bullets,
            diagram_png=diagram_png,
            assumptions=assumptions,
            dependencies=dependencies,
            risks=risks,
        )

        filename = f"{request.project_name.replace(' ', '_')}_ISBD.pptx"
        return StreamingResponse(
            pptx_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"ISBD generation error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate ISBD slides: {str(e)}")


# Save SoW Draft endpoint
class SoWDraftSaveRequest(BaseModel):
    """Request for saving SoW draft to file"""
    project_name: str
    customer: str
    background: str
    assumptions: Optional[str] = None
    out_of_scope: Optional[str] = None
    sow_content: str
    timestamp: str
    version: str = "1.0"


@app.post("/save/sow-draft")
async def save_sow_draft(request: SoWDraftSaveRequest):
    """
    Save SoW draft to the drafts folder as JSON file.
    """
    try:
        # Get the drafts folder path (relative to backend folder)
        backend_dir = Path(__file__).parent
        project_root = backend_dir.parent
        drafts_dir = project_root / "drafts"
        
        # Create drafts folder if it doesn't exist
        drafts_dir.mkdir(exist_ok=True)
        
        # Sanitize project name for filename
        sanitized_name = request.project_name.replace(' ', '_').replace('/', '_').replace('\\', '_')
        filename = f"SoW-Draft-{sanitized_name}.json"
        filepath = drafts_dir / filename
        
        # Prepare draft data
        draft_data = {
            "project_name": request.project_name,
            "customer": request.customer,
            "background": request.background,
            "assumptions": request.assumptions,
            "out_of_scope": request.out_of_scope,
            "sow_content": request.sow_content,
            "timestamp": request.timestamp,
            "version": request.version
        }
        
        # Write to file
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(draft_data, f, indent=2, ensure_ascii=False)
        
        logger.info(f"SoW draft saved: {filepath}")
        
        return {
            "success": True,
            "filename": filename,
            "path": str(filepath),
            "message": f"SoW draft saved successfully to drafts/{filename}"
        }
        
    except Exception as e:
        logger.error(f"Error saving SoW draft: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to save SoW draft: {str(e)}")


# Load SoW Draft endpoint
@app.get("/load/sow-draft/{project_name}")
async def load_sow_draft(project_name: str):
    """
    Load existing SoW draft from the drafts folder by project name.
    """
    try:
        # Get the drafts folder path
        backend_dir = Path(__file__).parent
        project_root = backend_dir.parent
        drafts_dir = project_root / "drafts"
        
        # Sanitize project name for filename
        sanitized_name = project_name.replace(' ', '_').replace('/', '_').replace('\\', '_')
        filename = f"SoW-Draft-{sanitized_name}.json"
        filepath = drafts_dir / filename
        
        # Check if file exists
        if not filepath.exists():
            return {
                "success": False,
                "exists": False,
                "message": f"No SoW draft found for project: {project_name}"
            }
        
        # Read the file
        with open(filepath, 'r', encoding='utf-8') as f:
            draft_data = json.load(f)
        
        logger.info(f"SoW draft loaded: {filepath}")
        
        return {
            "success": True,
            "exists": True,
            "draft": draft_data,
            "filename": filename,
            "message": f"SoW draft loaded successfully from drafts/{filename}"
        }
        
    except json.JSONDecodeError as e:
        logger.error(f"Error parsing SoW draft JSON: {e}")
        raise HTTPException(status_code=500, detail=f"Invalid JSON in SoW draft file: {str(e)}")
    except Exception as e:
        logger.error(f"Error loading SoW draft: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to load SoW draft: {str(e)}")


# Configuration Management Endpoints

class ProviderConfig(BaseModel):
    provider: str = Field(..., description="Provider name: anthropic, openai, or google_gemini")
    api_key: str = Field(..., description="API key for the provider")
    enabled: bool = Field(default=True, description="Enable this provider")


@app.get("/config/providers")
async def get_provider_config():
    """
    Get current provider configuration status.
    """
    try:
        backend_dir = Path(__file__).parent
        config_path = backend_dir / "config.json"
        
        if not config_path.exists():
            raise HTTPException(status_code=404, detail="Config file not found")
        
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)
        
        # Extract provider status
        providers = {}
        for provider_key in ['anthropic', 'openai', 'google_gemini']:
            if provider_key in config.get('models', {}):
                provider_data = config['models'][provider_key]
                providers[provider_key] = {
                    'enabled': provider_data.get('enabled', False),
                    'provider': provider_data.get('provider', ''),
                    'has_api_key': bool(os.getenv(provider_data.get('auth', {}).get('api_key_env', '')))
                }
        
        return {
            "success": True,
            "providers": providers,
            "default_provider": config.get('settings', {}).get('default_provider', 'sap_ai_core')
        }
        
    except Exception as e:
        logger.error(f"Error reading config: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to read configuration: {str(e)}")


@app.post("/config/update-provider")
async def update_provider_config(config_data: ProviderConfig):
    """
    Update provider configuration with API key.
    Updates both config.json and .env file.
    """
    try:
        backend_dir = Path(__file__).parent
        config_path = backend_dir / "config.json"
        env_path = backend_dir / ".env"
        
        # Validate provider
        valid_providers = {
            'anthropic': 'ANTHROPIC_API_KEY',
            'openai': 'OPENAI_API_KEY',
            'google_gemini': 'GOOGLE_API_KEY'
        }
        
        if config_data.provider not in valid_providers:
            raise HTTPException(status_code=400, detail=f"Invalid provider. Must be one of: {', '.join(valid_providers.keys())}")
        
        # Update config.json
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)
        
        provider_key = config_data.provider
        if provider_key not in config.get('models', {}):
            raise HTTPException(status_code=404, detail=f"Provider {provider_key} not found in config")
        
        # Enable the provider
        config['models'][provider_key]['enabled'] = config_data.enabled
        
        # If this is being enabled, also enable the first model
        if config_data.enabled and config['models'][provider_key].get('models'):
            config['models'][provider_key]['models'][0]['enabled'] = True
        
        # Save updated config
        with open(config_path, 'w', encoding='utf-8') as f:
            json.dump(config, f, indent=2)
        
        # Update or create .env file
        env_key = valid_providers[config_data.provider]
        env_lines = []
        key_found = False
        
        if env_path.exists():
            with open(env_path, 'r', encoding='utf-8') as f:
                env_lines = f.readlines()
            
            # Update existing key
            for i, line in enumerate(env_lines):
                if line.strip().startswith(f"{env_key}="):
                    env_lines[i] = f"{env_key}={config_data.api_key}\n"
                    key_found = True
                    break
        
        # Add key if not found
        if not key_found:
            env_lines.append(f"{env_key}={config_data.api_key}\n")
        
        # Write .env file
        with open(env_path, 'w', encoding='utf-8') as f:
            f.writelines(env_lines)
        
        # Update environment variable in current process
        os.environ[env_key] = config_data.api_key
        
        logger.info(f"Updated configuration for provider: {config_data.provider}")
        
        return {
            "success": True,
            "message": f"Successfully updated {config_data.provider} configuration",
            "provider": config_data.provider,
            "enabled": config_data.enabled,
            "restart_required": True
        }
        
    except json.JSONDecodeError as e:
        logger.error(f"Error parsing config JSON: {e}")
        raise HTTPException(status_code=500, detail=f"Invalid JSON in config file: {str(e)}")
    except Exception as e:
        logger.error(f"Error updating config: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to update configuration: {str(e)}")


# Public Holidays endpoint
HOLIDAYS_REGION_FILES = {
    "vic_australia": "holidays_vic_australia.json",
    "india": "holidays_india.json",
}


class HolidaysParseRequest(BaseModel):
    """Request to parse pasted holiday text into a JSON file for a region"""
    region: Literal["vic_australia", "india"]
    raw_text: str


@app.post("/holidays/parse")
async def parse_holidays(request: HolidaysParseRequest):
    """
    Use the active AI provider to parse pasted holiday text (copied from a
    webpage or spreadsheet) into structured holiday data for the current
    year only, then overwrite the region's JSON file under backend/data/.
    """
    try:
        if not request.raw_text.strip():
            raise HTTPException(status_code=400, detail="No holiday text provided")

        current_year = datetime.now().year
        region_label = "Victoria, Australia" if request.region == "vic_australia" else "India"

        prompt = f"""You are given raw text copied from a webpage or spreadsheet listing public holidays for {region_label}.

Extract only the holidays that fall in the year {current_year}.

Raw text:
---
{request.raw_text}
---

Return ONLY a JSON array (no markdown, no explanation) where each item has this shape:
{{"date": "YYYY-MM-DD", "name": "Holiday Name"}}

Sort the array by date ascending. If a holiday has no clear date or is not in {current_year}, omit it."""

        content = generate_with_active_provider(prompt=prompt, max_tokens=4000)
        content = content.strip()

        start_idx = content.find('[')
        end_idx = content.rfind(']') + 1
        if start_idx == -1 or end_idx <= start_idx:
            raise HTTPException(status_code=500, detail="AI response did not contain a JSON array")

        json_str = content[start_idx:end_idx]
        holidays = json.loads(json_str)

        backend_dir = Path(__file__).parent
        data_dir = backend_dir / "data"
        data_dir.mkdir(exist_ok=True)

        filepath = data_dir / HOLIDAYS_REGION_FILES[request.region]
        holiday_data = {
            "region": region_label,
            "year": current_year,
            "holidays": holidays,
            "updated_at": datetime.now().isoformat()
        }

        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(holiday_data, f, indent=2, ensure_ascii=False)

        logger.info(f"Holidays saved for {region_label}: {filepath} ({len(holidays)} entries)")

        return {
            "success": True,
            "region": request.region,
            "year": current_year,
            "count": len(holidays),
            "holidays": holidays
        }

    except json.JSONDecodeError as e:
        logger.error(f"Error parsing holidays JSON from AI response: {e}")
        raise HTTPException(status_code=500, detail=f"Could not parse AI response as JSON: {str(e)}")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Holidays parse error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to parse holidays: {str(e)}")


@app.get("/holidays/{region}")
async def get_holidays(region: Literal["vic_australia", "india"]):
    """Load the saved holidays JSON for a region, if it exists."""
    backend_dir = Path(__file__).parent
    filepath = backend_dir / "data" / HOLIDAYS_REGION_FILES[region]

    if not filepath.exists():
        return {"success": True, "region": region, "year": datetime.now().year, "holidays": []}

    with open(filepath, 'r', encoding='utf-8') as f:
        data = json.load(f)

    return {"success": True, **data}


if __name__ == "__main__":
    import uvicorn
    
    # Run the server
    uvicorn.run(
        "main:app",
        host="0.0.0.0",
        port=8000,
        reload=True,
        log_level="info"
    )
